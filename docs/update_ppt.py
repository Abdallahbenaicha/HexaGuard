"""Update PPT with corrected team names, email, and innovation highlights."""
import os, copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

PPTX = os.path.join(os.path.dirname(__file__),
                    "reference_files",
                    "Canevas-label-Projet-Innovant-2025-final.pptx")


def get_ref_rpr(tf):
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn("a:rPr"))
            return rPr
    return None

def get_ref_ppr(tf):
    if tf.paragraphs:
        return tf.paragraphs[0]._p.find(qn("a:pPr"))
    return None

def set_shape_text(shape, lines):
    tf = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody
    ref_rPr = get_ref_rpr(tf)
    ref_pPr = get_ref_ppr(tf)
    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)
    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if ref_pPr is not None:
            new_p.append(copy.deepcopy(ref_pPr))
        new_r = etree.SubElement(new_p, qn("a:r"))
        new_rPr = copy.deepcopy(ref_rPr) if ref_rPr is not None else etree.Element(qn("a:rPr"))
        new_r.insert(0, new_rPr)
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line
    print(f"  OK  '{shape.name}' -> {len(lines)} paragraphe(s)")

def find_shape(slide, key):
    for s in slide.shapes:
        if s.has_text_frame and key.lower() in s.text_frame.text.lower():
            return s
    return None

prs = Presentation(PPTX)
slides = prs.slides

# ── Slide 3 — CONTACT (email + noms corriges) ────────────────────────────
print("Slide 3 — Contact")
s = find_shape(slides[2], "Equipe SecurAx")
if not s:
    s = find_shape(slides[2], "Nom(s) du")
if s:
    set_shape_text(s, [
        "Nom(s) du responsable(s) : Benaicha Abdallah  |  Djenidi Ilies",
        "Fonctions : Architecte General & Expert Securite  |  Dev IA/Full-Stack",
        "Email : securaxdz@gmail.com",
        "Numero de telephone : [A completer]",
        "Site web (optionnel) : [En cours de deploiement]",
        "Reseaux sociaux : GitHub — SecurAx Project",
    ])

# ── Slide 5 — EQUIPE (noms reels) ─────────────────────────────────────────
print("Slide 5 — Equipe")
# Membres
s = find_shape(slides[4], "Prenom NOM")
if not s:
    s = find_shape(slides[4], "Chef de projet")
if s:
    set_shape_text(s, [
        "* Benaicha Abdallah — Architecte General, Scanners (Web/Reseau/SAST/DAST/Dep)",
        "                      & Developpement Server Internal Audit",
        "* Djenidi Ilies     — Agent IA ARIA (Gemini/Ollama), Frontend React/Vite/Tailwind,",
        "                      Backend Flask, API REST & Integration systeme",
    ])

# Competences
s = find_shape(slides[4], "Competences")
if s:
    set_shape_text(s, [
        "Benaicha Abdallah : Python, Nmap, SAST/DAST, CVSS v3.1, OWASP, Server hardening,",
        "  architecture systeme multi-modules, moteur de risque (FAIR + CISA KEV).",
        "Djenidi Ilies : React 18/Vite/Tailwind, Flask, SQLite/PostgreSQL, Google Gemini API,",
        "  Ollama (LLM local), ReportLab PDF, Flask-Login, 2FA TOTP, CORS/CSP securise.",
    ])

# ── Slide 7 — SOLUTION (aspect innovant integre) ──────────────────────────
print("Slide 7 — Solution (innovation)")
s = find_shape(slides[6], "Moteur de risque")
if s:
    set_shape_text(s, [
        "Moteur de risque : CVSS v3.1 + OWASP Risk Rating + FAIR + CISA KEV (24h)",
        "Rapports PDF bilingues Francais/Arabe — generes automatiquement en secondes",
        "Authentification renforcee : 2FA TOTP (RFC 6238), sessions securisees (CSP/HSTS)",
        "Dashboard temps reel : historique complet, export CSV / JSON / Markdown",
        "",
        "ASPECTS INNOVANTS :",
        "  SOUVERAINETE NUMERIQUE : solution algerienne, donnees hebergees localement,",
        "    independance des plateformes etrangeres (AWS, GCP, Azure non obligatoires).",
        "  COUT REDUIT           : 10x moins cher que Nessus/Qualys — accessible aux PME.",
        "  PLATEFORME UNIFIEE    : 7 outils de scan reunis en 1 seule interface coherente,",
        "    aucun besoin d'acheter/configurer des outils separes (Nmap + Burp + Snyk...).",
        "  STANDARDS MONDIAUX    : OWASP, NIST CVSS v3.1, ISO/IEC 27001, CISA KEV, GDPR.",
    ])

# ── Slide 10 — VALEUR AJOUTEE (reinforce aspect innovant) ─────────────────
print("Slide 10 — Valeur ajoutee")
s = find_shape(slides[9], "Ce que SecurAx")
if not s:
    s = find_shape(slides[9], "GAIN DE TEMPS")
if s:
    set_shape_text(s, [
        "Ce que SecurAx apporte — valeur reelle pour les clients :",
        "",
        "GAIN DE TEMPS        : Audit complet en 5 min vs. 5 jours manuellement (-80%)",
        "ECONOMIES            : 10x moins cher que Nessus / Qualys / Rapid7",
        "SOUVERAINETE NUMERIQUE : Solution made-in-Algeria, hebergement local possible,",
        "                         donnees sensibles restent sur le territoire national.",
        "PLATEFORME UNIFIEE   : 7 modules de scan + IA + rapports — tout en 1 outil,",
        "                       fini la fragmentation entre Nmap, Burp, Snyk, etc.",
        "INNOVATION IA (ARIA) : Analyse conversationnelle + correctifs proposes en FR/AR",
        "STANDARDS MONDIAUX   : OWASP, CVSS v3.1, ISO 27001, CISA KEV, GDPR — integres",
        "RAPPORTS PRÊTS       : PDF, CSV, JSON, Markdown — pour auditeurs et direction",
        "ACCESSIBILITE        : Interface bilingue AR/FR, adaptee au marche algerien",
    ])

# ── Slide 11 — CONCURRENCE (souverainete comme differenciation) ────────────
print("Slide 11 — Concurrence (souverainete)")
s = find_shape(slides[10], "Concurrents directs")
if s:
    set_shape_text(s, [
        "Concurrents directs  : Nessus (Tenable), Qualys, Rapid7, Burp Suite Pro",
        "Concurrents indirects: OWASP ZAP (gratuit/manuel), Snyk (dependances seul)",
        "",
        "Analyse SWOT :",
        "  FORCES      : Prix accessible, IA integree, FR/AR, tout-en-un, LOCAL & souverain",
        "  FAIBLESSES  : Notoriete limitee au demarrage, ressources initiales reduites",
        "  OPPORTUNITES: Marche algerien/africain sous-equipe, hausse cyberattaques +38%",
        "  MENACES     : Geants (Tenable, IBM), nouvelles startups IA securite mondiales",
        "",
        "Differenciation cle — pourquoi SecurAx est unique :",
        "  (1) SOUVERAINETE NUMERIQUE : hebergement 100% algerien possible, pas de cloud etranger",
        "  (2) COUT ACCESSIBLE        : a partir de 29 $/mois vs. 3 000 $/an pour Nessus",
        "  (3) PLATEFORME UNIFIEE     : 7 scanners + ARIA IA dans 1 seul outil coherent",
        "  (4) STANDARDS INTERNATIONAUX: CVSS v3.1, OWASP, CISA KEV — meme niveau que leaders",
        "  (5) BILINGUE FR/AR         : seule solution adaptee au marche algerien et africain",
    ])

# ── SAVE ──────────────────────────────────────────────────────────────────
prs.save(PPTX)
print()
print("SAVED ->", PPTX)
