"""
Fill Canevas-label-Projet-Innovant-2025-final.pptx with SecurAx project data.
Run from any directory: python fill_ppt.py
"""
import os
import copy
from pptx import Presentation
from pptx.oxml.ns import qn
from lxml import etree

PPTX = os.path.join(os.path.dirname(__file__),
                    "reference_files",
                    "Canevas-label-Projet-Innovant-2025-final.pptx")


def get_ref_rpr(tf):
    for para in tf.paragraphs:
        for run in para.runs:
            rPr = run._r.find(qn("a:rPr"))
            return rPr
    return None


def get_ref_ppr(tf):
    if tf.paragraphs:
        return tf.paragraphs[0]._p.find(qn("a:pPr"))
    return None


def set_shape_text(shape, lines, font_size_pt=None, bold=None, color_hex=None):
    tf = shape.text_frame
    tf.word_wrap = True
    txBody = tf._txBody

    ref_rPr = get_ref_rpr(tf)
    ref_pPr = get_ref_ppr(tf)

    for p in txBody.findall(qn("a:p")):
        txBody.remove(p)

    for line in lines:
        new_p = etree.SubElement(txBody, qn("a:p"))
        if ref_pPr is not None:
            new_p.append(copy.deepcopy(ref_pPr))

        new_r = etree.SubElement(new_p, qn("a:r"))

        if ref_rPr is not None:
            new_rPr = copy.deepcopy(ref_rPr)
        else:
            new_rPr = etree.Element(qn("a:rPr"))
            new_rPr.set("lang", "fr-FR")
            new_rPr.set("dirty", "0")

        if font_size_pt is not None:
            new_rPr.set("sz", str(int(font_size_pt * 100)))
        if bold is not None:
            new_rPr.set("b", "1" if bold else "0")
        if color_hex is not None:
            for sf in new_rPr.findall(qn("a:solidFill")):
                new_rPr.remove(sf)
            solidFill = etree.SubElement(new_rPr, qn("a:solidFill"))
            srgbClr = etree.SubElement(solidFill, qn("a:srgbClr"))
            srgbClr.set("val", color_hex)

        new_r.insert(0, new_rPr)
        new_t = etree.SubElement(new_r, qn("a:t"))
        new_t.text = line

    print(f"  OK  shape '{shape.name}' -> {len(lines)} paragraph(s)")


def find_shape(slide, partial_text):
    for s in slide.shapes:
        if s.has_text_frame and partial_text.lower() in s.text_frame.text.lower():
            return s
    return None


def fs(slide, key):
    s = find_shape(slide, key)
    if s is None:
        print(f"  WARN  shape not found for key: {key!r}")
    return s


prs = Presentation(PPTX)
slides = prs.slides

# ── Slide 2 — TITRE DU PROJET ─────────────────────────────────────────────
print("Slide 2 — Titre")
sl = slides[1]
s = fs(sl, "Nom principal")
if s:
    set_shape_text(s, ["SecurAx"])

s = fs(sl, "logo officiel")
if s:
    set_shape_text(s, ["[ Logo SecurAx — a inserer ]"])

s = fs(sl, "phrase courte")
if s:
    set_shape_text(s, [
        "La cybersecurite intelligente, automatisee et accessible —",
        "analysez, protegez et repondez en quelques minutes.",
    ])

# ── Slide 3 — INFORMATIONS DE CONTACT ────────────────────────────────────
print("Slide 3 — Contact")
sl = slides[2]
s = fs(sl, "Nom(s) du")
if s:
    set_shape_text(s, [
        "Nom(s) du responsable(s) : Equipe SecurAx",
        "Fonction : Ingenieurs en Cybersecurite & Developpeurs Full-Stack",
        "Email : innovation.team.dz@gmail.com",
        "Numero de telephone : [A completer]",
        "Site web (optionnel) : [En cours de deploiement]",
        "Reseaux sociaux (optionnel) : GitHub — SecurAx Project",
    ])

# ── Slide 4 — RESUME DU PROJET ────────────────────────────────────────────
print("Slide 4 — Resume")
sl = slides[3]
s = fs(sl, "PRÉSENTEZ")
if not s:
    s = fs(sl, "PRESENTEZ")
if s:
    set_shape_text(s, [
        "SecurAx est une plateforme algerienne de cybersecurite qui automatise",
        "l'evaluation des vulnerabilites grace a l'intelligence artificielle.",
        "",
        "PROBLEMATIQUE : 60% des PME algeriennes n'ont jamais realise d'audit",
        "de securite. Les solutions pro coutent 3 000 $ - 50 000 $/an.",
        "",
        "SOLUTION : 7 modules de scan (Web, Reseau, SAST, DAST, Dependances,",
        "Serveur, Apache) + agent IA ARIA (Gemini/Ollama) + rapports PDF auto.",
        "",
        "IMPACT : -80% de temps d'audit — rapports complets en 5 minutes.",
        "Cout 10x inferieur aux solutions commerciales existantes.",
    ])

# ── Slide 5 — L'EQUIPE ───────────────────────────────────────────────────
print("Slide 5 — Equipe")
sl = slides[4]
s = fs(sl, "PRÉSENTATION DES MEMBRES")
if not s:
    s = fs(sl, "PRESENTATION DES MEMBRES")
if s:
    set_shape_text(s, [
        "* [Prenom NOM] — Chef de projet, Architecte IA & Backend Python/Flask",
        "* [Prenom NOM] — Developpeur Frontend React / UI-UX / Tailwind CSS",
        "* [Prenom NOM] — Expert Cybersecurite & Tests de penetration",
        "* [Prenom NOM] — Developpeur Backend & Base de donnees / DevOps",
    ])

s = fs(sl, "COMPÉTENCES")
if not s:
    s = fs(sl, "COMPETENCES")
if s:
    set_shape_text(s, [
        "Competences : Python, Flask, React/Vite/Tailwind, SQLite/PostgreSQL,",
        "IA generative (Gemini, Ollama), Securite (CVSS v3.1, OWASP, NMAP, SAST).",
        "Roles : Architecture, dev full-stack, integration IA, rapports PDF (ReportLab),",
        "authentification 2FA TOTP, moteur de risque multicritere (CVSS+FAIR+KEV).",
    ])

# ── Slide 6 — PROBLEMATIQUE ───────────────────────────────────────────────
print("Slide 6 — Problematique")
sl = slides[5]
s = fs(sl, "QUEL EST LE PROBLÈME")
if not s:
    s = fs(sl, "QUEL EST LE PROBLEME")
if s:
    set_shape_text(s, [
        "Probleme : Les entreprises algeriennes et africaines manquent d'outils",
        "de cybersecurite accessibles — 60% des PME n'ont jamais realise d'audit.",
    ])

s = fs(sl, "DONNÉES CHIFFRÉES")
if not s:
    s = fs(sl, "DONNEES CHIFFREES")
if not s:
    s = fs(sl, "CHIFFR")
if s:
    set_shape_text(s, [
        "> Cout moyen d'une violation de donnees : 4,45 M$ (IBM 2023)",
        "> +38% de cyberattaques mondiales en 2022 (Check Point Research)",
        "> Delai moyen de detection d'une intrusion sans outil : 207 jours",
        "> Seulement 12% des startups algeriennes ont un plan de securite",
    ])

s = fs(sl, "VISUELS")
if s:
    set_shape_text(s, [
        "> Solutions concurrentes (Nessus, Qualys, Rapid7) : 3 000 $ - 50 000 $/an",
        "  -> Inaccessibles pour startups, developpeurs independants et PME locales.",
        "> Absence de plateforme FR/AR dediee au marche algerien et africain.",
    ])

# ── Slide 7 — LA SOLUTION ─────────────────────────────────────────────────
print("Slide 7 — Solution")
sl = slides[6]
s = fs(sl, "DESCRIPTION CONCISE")
if s:
    set_shape_text(s, [
        "SecurAx propose 7 modules de scan integres dans une interface unifiee :",
        "",
        "(1) Web Scanner        — Detection OWASP Top 10 (XSS, SQLi, CSRF...)",
        "(2) Network Scanner    — Reconnaissance reseau Nmap, ports, OS fingerprinting",
        "(3) SAST               — Analyse statique du code source (ZIP) multi-langages",
        "(4) DAST               — Tests dynamiques sur applications web en production",
        "(5) Dependency Scanner — Audit CVE/NVD des librairies (npm, pip, Maven)",
        "(6) Server/Apache Audit— Configuration systeme, hardening, directives Apache",
        "(7) ARIA IA            — Agent conversationnel (Gemini/Ollama) pour remediation",
    ])

s = fs(sl, "ORIGINALITÉ")
if not s:
    s = fs(sl, "ORIGINALITE")
if not s:
    s = fs(sl, "SOLUTION, EN METTANT")
if s:
    set_shape_text(s, [
        "Moteur de risque : CVSS v3.1 + OWASP Risk Rating + FAIR + CISA KEV",
        "Rapports PDF professionnels bilingues (Francais/Arabe) — auto-generes",
        "Authentification renforcee : 2FA TOTP (RFC 6238), sessions securisees",
        "Dashboard : historique des scans, stats, export CSV/JSON/Markdown",
    ])

# ── Slide 8 — PROTOTYPE ───────────────────────────────────────────────────
print("Slide 8 — Prototype")
sl = slides[7]
s = fs(sl, "METTEZ ICI LE LIEN")
if s:
    set_shape_text(s, [
        "Prototype fonctionnel : SecurAx v1.0 (Beta)",
        "",
        "Stack technique :",
        "  Backend  : Flask (Python 3.11), SQLite/PostgreSQL, Gunicorn + Docker",
        "  Frontend : React 18 + Vite + Tailwind CSS (SPA)",
        "  IA       : Google Gemini Pro API / Ollama (fallback local hors-ligne)",
        "  Securite : Flask-WTF CSRF, Flask-Login, Flask-Limiter, CORS securise",
    ])

s = fs(sl, "VIDÉO EXPLICATIVE")
if not s:
    s = fs(sl, "VIDEO EXPLICATIVE")
if not s:
    s = fs(sl, "APK")
if s:
    set_shape_text(s, [
        "Demo video   : [ lien a inserer — 4 min max ]",
        "Code source  : [ GitHub — SecurAx — lien a inserer ]",
        "Application  : [ https://securax.dz — deploiement en cours ]",
        "Note         : Application web SaaS (pas d'APK Android/iOS prevu en phase 1)",
    ])

# ── Slide 9 — PROPRIETE INTELLECTUELLE ───────────────────────────────────
print("Slide 9 — PI")
sl = slides[8]
s = fs(sl, "STATUT DE LA PROPRI")
if s:
    set_shape_text(s, [
        "Statut de la propriete intellectuelle :",
        "",
        "Brevets deposes ou en cours :",
        "  -> Brevet algerien (INAPI) prevu : moteur de risque multicritere SecurAx",
        "",
        "Marques ou designs enregistres :",
        "  -> Marque 'SecurAx' — enregistrement INAPI en cours",
        "",
        "Droits d'auteur et licences :",
        "  -> Code source sous licence proprietaire (droits reserves — equipe SecurAx)",
        "  -> Composants open-source : licences MIT, Apache 2.0, BSD (conformite verifiee)",
    ])

# ── Slide 10 — VALEUR AJOUTEE ────────────────────────────────────────────
print("Slide 10 — Valeur ajoutee")
sl = slides[9]
s = fs(sl, "CE QUE VOUS APPORTEZ")
if s:
    set_shape_text(s, [
        "Ce que SecurAx apporte aux clients :",
        "",
        "GAIN DE TEMPS      : Audit complet en 5 min vs. 5 jours manuellement (-80%)",
        "ECONOMIES          : 10x moins cher que les solutions professionnelles",
        "INNOVATION IA      : ARIA analyse et propose des correctifs en langage naturel",
        "RAPPORTS CLES      : PDF, CSV, JSON, Markdown — prets pour direction et auditeurs",
        "LOCALISATION       : Plateforme bilingue FR/AR, adaptee au marche algerien",
        "CONFORMITE         : Alignee OWASP, CVSS v3.1, GDPR, ISO/IEC 27001, CISA KEV",
        "SCALABILITE        : SaaS multi-tenant, de l'independant a la grande entreprise",
    ])

# ── Slide 11 — ANALYSE CONCURRENTIELLE ───────────────────────────────────
print("Slide 11 — Concurrence")
sl = slides[10]
s = fs(sl, "ANALYSE DE LA CONCURRENCE")
if s:
    set_shape_text(s, [
        "Concurrents directs : Nessus (Tenable), Qualys, Rapid7 InsightVM, Burp Suite Pro",
        "Concurrents indirects : OWASP ZAP (gratuit/manuel), Snyk (dependances seulement)",
        "",
        "Analyse SWOT :",
        "  FORCES      : Prix accessible, IA integree, FR/AR, tout-en-un, solution locale",
        "  FAIBLESSES  : Notoriete limitee, ressources de demarrage reduites",
        "  OPPORTUNITES: Marche algerien/africain sous-equipe, +38% cyberattaques/an",
        "  MENACES     : Geants (Tenable, IBM), nouvelles startups IA de securite",
        "",
        "Differenciation cle :",
        "  Seule plateforme algerienne combinant 7 scans + ARIA IA + rapports FR/AR",
        "  + moteur de risque multicritere dans une offre accessible (a partir de 29 $/mois).",
    ])

# ── Slide 12 — MARCHE CIBLE ───────────────────────────────────────────────
print("Slide 12 — Marche")
sl = slides[11]
s = fs(sl, "CIBLE PRINCIPALE")
if s:
    set_shape_text(s, [
        "Cible principale :",
        "  (1) Developpeurs & equipes techniques (startups, freelances) — abonnement mensuel",
        "  (2) PME et TPE algeriennes et africaines — audit ponctuel ou contrat annuel",
        "  (3) Entreprises et etablissements publics — licence entreprise + support dedie",
        "  (4) Agences cybersecurite & pentesters — integration dans leur workflow",
        "",
        "Taille du marche :",
        "  Marche mondial cybersecurite : 202 Md$ (2023) -> 500 Md$ (2030)",
        "  Marche Afrique du Nord       : croissance +18%/an (Statista 2023)",
        "  Cible immediate              : ~50 000 developpeurs & PME en Algerie",
    ])

# ── Slide 13 — BUSINESS MODEL ────────────────────────────────────────────
print("Slide 13 — Business Model")
sl = slides[12]
s = fs(sl, "GÉNÉRER DES REVENUS")
if not s:
    s = fs(sl, "GENERER DES REVENUS")
if s:
    set_shape_text(s, [
        "Comment SecurAx genere des revenus :",
        "",
        "(1) Freemium SaaS :",
        "    Gratuit  : 2 scans/mois, rapport basique",
        "    Pro      : 29 $/mois — scans illimites, tous modules, ARIA IA, PDF complet",
        "    Enterprise : sur devis — multi-utilisateurs, API, CI/CD, SLA garanti",
        "",
        "(2) Audits ponctuels a la carte : 99 $ - 499 $ selon le type de scan",
        "",
        "(3) API SecurAx : Integration DevSecOps (GitHub Actions, GitLab CI) — pay-per-scan",
        "",
        "(4) Formations & Certifications : e-learning en cybersecurite offensive/defensive",
        "",
        "Projection An 1 : 200 clients Pro -> CA ~70 K$/an",
        "Projection An 3 : 2 000 clients (Pro + Enterprise) -> CA ~800 K$/an",
    ])

# ── Slide 14 — ROADMAP ────────────────────────────────────────────────────
print("Slide 14 — Roadmap")
sl = slides[13]
s = fs(sl, "ÉTAPES CLÉS")
if not s:
    s = fs(sl, "ETAPES CLES")
if s:
    set_shape_text(s, [
        "Roadmap et strategie de developpement :",
        "",
        "[FAIT 2025] Developpement v1.0 :",
        "  7 modules de scan, agent ARIA, rapports PDF, 2FA TOTP, dashboard React",
        "  Architecture Blueprint Flask + React SPA + SQLite (prete production)",
        "",
        "[6 MOIS — dec. 2025] Lancement Beta :",
        "  Deploiement cloud Algerie, 50 beta-testeurs, integration API GitHub Actions",
        "  Ressources : 2 serveurs VPS (2 000 $/an), 3 devs a temps partiel",
        "",
        "[1 AN — juin 2026] Lancement commercial :",
        "  200 clients Pro, 5 partenariats entreprises algeriennes, ISO 27001",
        "  Ressources : equipe 5 personnes, levee 150 K$ (subvention Algerie Startup)",
        "",
        "[3 ANS — 2028] Expansion regionale :",
        "  Marches : Tunisie, Maroc, Senegal, France — 2 000+ clients, API publique",
        "  Module IA avance (LLM fine-tune securite), CA > 800 K$, Serie A",
    ])

# ── SAVE ──────────────────────────────────────────────────────────────────
prs.save(PPTX)
print()
print("SAVED ->", PPTX)
