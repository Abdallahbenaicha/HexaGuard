"""
Fix text overflow / overlap in the SecurAx PPT.
Strategy:
  - Reduce font size in all content boxes to 14pt
  - Resize / reposition content boxes so they don't overlap each other
  - All boxes extended to near the slide bottom (26 cm)
"""
import os
from pptx import Presentation
from pptx.util import Pt, Emu
from pptx.oxml.ns import qn

PPTX = os.path.join(os.path.dirname(__file__),
                    "reference_files",
                    "Canevas-label-Projet-Innovant-2025-final.pptx")

CM = 360000   # 1 cm in EMU


def cm(v):
    return int(v * CM)


def set_font_size(shape, pt):
    """Set font size for every run inside a shape."""
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    for para in tf.paragraphs:
        # paragraph-level default run properties
        pPr = para._p.find(qn("a:pPr"))
        if pPr is not None:
            dRPr = pPr.find(qn("a:defRPr"))
            if dRPr is not None:
                dRPr.set("sz", str(int(pt * 100)))
        # run-level properties
        for run in para.runs:
            rPr = run._r.find(qn("a:rPr"))
            if rPr is not None:
                rPr.set("sz", str(int(pt * 100)))
            run.font.size = Pt(pt)


def resize(shape, top_cm=None, height_cm=None, left_cm=None, width_cm=None):
    if top_cm is not None:
        shape.top = cm(top_cm)
    if height_cm is not None:
        shape.height = cm(height_cm)
    if left_cm is not None:
        shape.left = cm(left_cm)
    if width_cm is not None:
        shape.width = cm(width_cm)


def find(slide, key):
    for s in slide.shapes:
        if s.has_text_frame and key.lower() in s.text_frame.text.lower():
            return s
    return None


prs = Presentation(PPTX)
slides = prs.slides
BOTTOM = 26.0   # safe bottom margin in cm


# ── Slide 2 — TITRE (short text, 32pt — just extend slightly) ────────────
sl = slides[1]
for key in ["SecurAx", "Logo SecurAx", "cybersecurite intelligente"]:
    s = find(sl, key)
    if s:
        top_cm = round(s.top / CM, 1)
        resize(s, height_cm=BOTTOM - top_cm)
        set_font_size(s, 22)
        print(f"Slide 2 resized '{key[:30]}'")

# ── Slide 3 — CONTACT (single box, 32pt → 16pt) ───────────────────────────
sl = slides[2]
s = find(sl, "Benaicha Abdallah")
if not s:
    s = find(sl, "Nom(s) du")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 18)
    print(f"Slide 3 resized contact box")

# ── Slide 4 — RESUME (single box) ─────────────────────────────────────────
sl = slides[3]
s = find(sl, "SecurAx est une plateforme")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 16)
    print(f"Slide 4 resized resume box")

# ── Slide 5 — EQUIPE (2 boxes — reposition + resize) ──────────────────────
sl = slides[4]
# Box 1: members (was top=6.3, h=1.2)
b1 = find(sl, "Benaicha Abdallah")
# Box 2: competences (was top=8.0, h=1.2)
b2 = find(sl, "Competences")
if b1 and b2:
    resize(b1, top_cm=6.3, height_cm=8.0, width_cm=44.0)
    set_font_size(b1, 16)
    resize(b2, top_cm=15.0, height_cm=BOTTOM - 15.0, width_cm=44.0)
    set_font_size(b2, 16)
    print(f"Slide 5 repositioned 2 boxes")

# ── Slide 6 — PROBLEMATIQUE (3 boxes — stack properly) ────────────────────
sl = slides[5]
b1 = find(sl, "Probleme : Les entreprises")
b2 = find(sl, "Cout moyen")
b3 = find(sl, "Solutions concurrentes")
if b1:
    resize(b1, top_cm=6.3, height_cm=3.5, width_cm=44.0)
    set_font_size(b1, 16)
if b2:
    resize(b2, top_cm=10.5, height_cm=5.0, width_cm=44.0)
    set_font_size(b2, 16)
if b3:
    resize(b3, top_cm=16.5, height_cm=BOTTOM - 16.5, width_cm=44.0)
    set_font_size(b3, 16)
print(f"Slide 6 repositioned 3 boxes")

# ── Slide 7 — SOLUTION (2 dense boxes — split slide) ──────────────────────
sl = slides[6]
b1 = find(sl, "SecurAx propose 7 modules")
b2 = find(sl, "Moteur de risque")
if b1:
    resize(b1, top_cm=6.3, height_cm=10.0, width_cm=44.0)
    set_font_size(b1, 14)
if b2:
    resize(b2, top_cm=17.0, height_cm=BOTTOM - 17.0, width_cm=44.0)
    set_font_size(b2, 14)
print(f"Slide 7 repositioned 2 boxes")

# ── Slide 8 — PROTOTYPE (2 boxes) ─────────────────────────────────────────
sl = slides[7]
b1 = find(sl, "Prototype fonctionnel")
b2 = find(sl, "Demo video")
if b1:
    resize(b1, top_cm=6.3, height_cm=9.0, width_cm=44.0)
    set_font_size(b1, 16)
if b2:
    resize(b2, top_cm=16.0, height_cm=BOTTOM - 16.0, width_cm=44.0)
    set_font_size(b2, 16)
print(f"Slide 8 repositioned 2 boxes")

# ── Slide 9 — PI (single box) ─────────────────────────────────────────────
sl = slides[8]
s = find(sl, "propriete intellectuelle")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 16)
    print(f"Slide 9 resized")

# ── Slide 10 — VALEUR AJOUTEE ─────────────────────────────────────────────
sl = slides[9]
s = find(sl, "Ce que SecurAx apporte")
if not s:
    s = find(sl, "GAIN DE TEMPS")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 15)
    print(f"Slide 10 resized")

# ── Slide 11 — CONCURRENCE ────────────────────────────────────────────────
sl = slides[10]
s = find(sl, "Concurrents directs")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 14)
    print(f"Slide 11 resized")

# ── Slide 12 — MARCHE ─────────────────────────────────────────────────────
sl = slides[11]
s = find(sl, "Cible principale")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 16)
    print(f"Slide 12 resized")

# ── Slide 13 — BUSINESS MODEL ─────────────────────────────────────────────
sl = slides[12]
s = find(sl, "Comment SecurAx genere")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 15)
    print(f"Slide 13 resized")

# ── Slide 14 — ROADMAP ────────────────────────────────────────────────────
sl = slides[13]
s = find(sl, "Roadmap et strategie")
if s:
    top_cm = round(s.top / CM, 1)
    resize(s, height_cm=BOTTOM - top_cm, width_cm=44.0)
    set_font_size(s, 14)
    print(f"Slide 14 resized")

# ── SAVE ──────────────────────────────────────────────────────────────────
prs.save(PPTX)
print()
print("SAVED ->", PPTX)
